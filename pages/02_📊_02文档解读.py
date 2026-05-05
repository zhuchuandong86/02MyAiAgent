import streamlit as st
import os
import hashlib 
from datetime import datetime

import core.paths
from core.settings import settings
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document

from core.parsers.document_engine import convert_pdf_to_images
from modules.multi_compare.main import process_single_page, generate_final_summary
from modules.multi_compare.main_compare import generate_compare_summary
from modules.multi_compare.main_trend import generate_trend_summary
from modules.multi_compare.renderers.my_html_renderer import export_to_html
from modules.multi_compare.renderers.excel_builder import export_tables_to_excel

# 【统一引入 RAG 与提示词管家】
from modules.multi_compare.template_service import (
    get_available_templates, get_style_templates, get_embeddings, 
    TEMPLATE_MD_DIR, TEMPLATE_DB_DIR
)
from modules.multi_compare.ui_prompts import (
    UI_COGNITIVE_SINGLE, UI_COGNITIVE_COMPARE, UI_COGNITIVE_TREND, 
    UI_CHART_MERMAID, GET_USER_PRIORITY, GET_STYLE_FUSION
)

OUTPUT_DIR = str(core.paths.GLOBAL_DATA_DIR)
UPLOAD_DIR = str(core.paths.UPLOAD_DIR)
TEMP_IMG_DIR = os.path.join(OUTPUT_DIR, "pdf_temp_images")
MD_CACHE_DIR = os.path.join(OUTPUT_DIR, "md_cache") 

os.makedirs(TEMP_IMG_DIR, exist_ok=True)
os.makedirs(MD_CACHE_DIR, exist_ok=True)

st.set_page_config(page_title="AI 材料深度解读工作台", page_icon="📊", layout="wide", initial_sidebar_state="expanded")

# =========================================================
# 🌟 全局初始化 Session State (为三个 Tab 准备独立的记忆容器)
# =========================================================
# Tab1: 单文档
if "tab1_report" not in st.session_state: st.session_state.tab1_report = None
if "tab1_docs" not in st.session_state: st.session_state.tab1_docs = ""
if "tab1_history" not in st.session_state: st.session_state.tab1_history = []
if "tab1_docx_path" not in st.session_state: st.session_state.tab1_docx_path = None
if "tab1_base_name" not in st.session_state: st.session_state.tab1_base_name = "研判报告"
if "tab1_report_title" not in st.session_state: st.session_state.tab1_report_title = "研判报告"

# Tab2: 横向竞品
if "tab2_report" not in st.session_state: st.session_state.tab2_report = None
if "tab2_docs" not in st.session_state: st.session_state.tab2_docs = ""
if "tab2_history" not in st.session_state: st.session_state.tab2_history = []
if "tab2_prefix" not in st.session_state: st.session_state.tab2_prefix = "竞品横评"

# Tab3: 纵向趋势
if "tab3_report" not in st.session_state: st.session_state.tab3_report = None
if "tab3_docs" not in st.session_state: st.session_state.tab3_docs = ""
if "tab3_history" not in st.session_state: st.session_state.tab3_history = []
if "tab3_prefix" not in st.session_state: st.session_state.tab3_prefix = "演进趋势"

st.markdown("""
    <style>
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    .stMarkdown p { font-size: 16px !important; line-height: 1.7 !important; margin-bottom: 12px !important; }
    .stMarkdown li { margin-bottom: 6px !important; }
    .stMarkdown li > p { margin-bottom: 0px !important; margin-top: 0px !important; }
    .stMarkdown ul, .stMarkdown ol { margin-bottom: 20px !important; padding-left: 28px !important; }
    .stMarkdown h2 { font-size: 22px !important; color: #1a202c !important; border-bottom: 2px solid #ebf4ff !important; padding-bottom: 8px !important; margin-top: 35px !important; margin-bottom: 16px !important; }
    .stMarkdown h3 { font-size: 18px !important; color: #2b6cb0 !important; margin-top: 24px !important; margin-bottom: 12px !important; font-weight: 600 !important; }
    .stMarkdown strong { color: #111827 !important; }
    </style>
    """, unsafe_allow_html=True)

with st.sidebar:
    max_pages = st.number_input("原文件最大解析页数 (防超载)", min_value=1, max_value=500, value=100)
    st.markdown("---")
    use_cache = st.checkbox("⚡ 启用极速解析缓存", value=True)

st.markdown("### 📊 AI 材料深度解读工作台")
st.markdown("基于多模态大模型，支持上传 PDF、图片、MD、Word、PPT 混合格式，进行单文档解析、多文件横评与纵向趋势推演。")

# ================= render_export_buttons 函数 =================
def render_export_buttons(summary_md, base_filename, report_type="报告", docx_path=None):
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    md_name = f"{base_filename}_{report_type}_{timestamp}.md"
    html_name = f"{base_filename}_{report_type}_{timestamp}.html"
    excel_name = f"{base_filename}_数据表_{timestamp}.xlsx"
    
    final_md_file = os.path.join(OUTPUT_DIR, md_name)
    final_html_file = os.path.join(OUTPUT_DIR, html_name)
    final_excel_file = os.path.join(OUTPUT_DIR, excel_name)
    
    with open(final_md_file, "w", encoding="utf-8") as f: f.write(summary_md)
    export_to_html(summary_md, final_html_file)
    has_excel = export_tables_to_excel(summary_md, final_excel_file)
    
    st.markdown("### 💾 导出报告")
    
    col_count = 2
    if has_excel: col_count += 1
    if docx_path and os.path.exists(docx_path): col_count += 1
        
    cols = st.columns(col_count)
    curr_col = 0
    
    with cols[curr_col]:
        with open(final_md_file, "r", encoding="utf-8") as f: st.download_button("⬇️ Markdown版", f, file_name=md_name)
    curr_col += 1
        
    with cols[curr_col]:
        with open(final_html_file, "r", encoding="utf-8") as f: st.download_button("🌐 HTML网页版", f, file_name=html_name)
    curr_col += 1
        
    if has_excel:
        with cols[curr_col]:
            with open(final_excel_file, "rb") as f: st.download_button("📊 Excel数据表", f, file_name=excel_name, mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        curr_col += 1
            
    if docx_path and os.path.exists(docx_path):
        with cols[curr_col]:
            with open(docx_path, "rb") as f: st.download_button("📘 原生Word版 (DOCX)", f, file_name=f"{base_filename}_{report_type}_{timestamp}.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document", type="primary")

# ================= 强大的多模态文件解析 (无任何删减) =================
def parse_files_to_text_dict(uploaded_files, max_pages, ui_container, enable_cache):
    import os
    import hashlib
    from datetime import datetime
    result_dict = {}
    
    for file in uploaded_files:
        ext = os.path.splitext(file.name)[1].lower()
        base_name = os.path.splitext(file.name)[0]
        
        if ext == '.md':
            ui_container.info(f"📄 秒读 Markdown 文本: {file.name}")
            result_dict[base_name] = file.getvalue().decode("utf-8")
        else:
            file_bytes = file.getvalue()
            file_md5 = hashlib.md5(file_bytes).hexdigest()
            cache_filename = f"{file_md5}_limit_{max_pages}.md"
            cache_file_path = os.path.join(MD_CACHE_DIR, cache_filename)
            
            if enable_cache and os.path.exists(cache_file_path):
                ui_container.success(f"⚡ 物理指纹匹配成功！命中全系统级缓存: {file.name}")
                with open(cache_file_path, "r", encoding="utf-8") as f: 
                    result_dict[base_name] = f.read()
                continue
                
            all_content = ""
            
            if ext == '.docx':
                ui_container.info(f"📝 正在深度解剖 Word 文档: {file.name}")
                try:
                    import docx
                    import io
                    doc = docx.Document(io.BytesIO(file_bytes))
                    full_text = []
                    
                    for para in doc.paragraphs:
                        if para.text.strip(): full_text.append(para.text)
                    for table in doc.tables:
                        for row in table.rows:
                            full_text.append(" | ".join([cell.text.replace('\n', ' ') for cell in row.cells]))
                    
                    img_count = 0
                    for rel in doc.part.rels.values():
                        if "image" in rel.target_ref:
                            img_bytes = rel.target_part.blob
                            if len(img_bytes) > 15360: 
                                img_count += 1
                                ui_container.write(f"   🔍 发现有效 Word 插图 {img_count}，正在呼叫视觉引擎...")
                                img_path = os.path.join(TEMP_IMG_DIR, f"word_{file_md5[:8]}_{img_count}.png")
                                with open(img_path, "wb") as f: f.write(img_bytes)
                                vis_res = process_single_page(img_path, f"Word核心插图_{img_count}")
                                full_text.append(f"\n> 🖼️ **[文档插图/数据表 解析]**:\n{vis_res}\n")
                                
                    all_content = "\n".join(full_text)
                    ui_container.success(f"✅ {file.name} 文本与图片双路解析完成！")
                except ImportError:
                    ui_container.error("❌ 缺少 docx 库，请执行 `pip install python-docx`")
                    continue
                    
            elif ext == '.pptx':
                ui_container.info(f"📊 正在逐页深度解剖 PPT: {file.name}")
                try:
                    from pptx import Presentation
                    import io
                    ppt = Presentation(io.BytesIO(file_bytes))
                    full_text = []
                    
                    progress_bar = ui_container.progress(0)
                    total_slides = len(ppt.slides)
                    
                    for i, slide in enumerate(ppt.slides):
                        slide_text = []
                        img_count = 0
                        
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_text.append(shape.text.strip())
                                
                            if hasattr(shape, "image"):
                                img_bytes = shape.image.blob
                                if len(img_bytes) > 15360: 
                                    img_count += 1
                                    ui_container.write(f"   🔍 发现第 {i+1} 页核心插图，正在呼叫视觉引擎...")
                                    img_path = os.path.join(TEMP_IMG_DIR, f"ppt_{file_md5[:8]}_s{i+1}_i{img_count}.png")
                                    with open(img_path, "wb") as f: f.write(img_bytes)
                                    vis_res = process_single_page(img_path, f"第{i+1}页_插图{img_count}")
                                    slide_text.append(f"\n> 🖼️ **[本页核心插图/架构图 解析]**:\n{vis_res}\n")
                                    
                        full_text.append(f"\n--- 📑 第 {i+1} 页幻灯片 ---\n" + "\n".join(slide_text))
                        progress_bar.progress((i + 1) / total_slides)
                        
                    progress_bar.empty()
                    all_content = "\n".join(full_text)
                    ui_container.success(f"✅ {file.name} PPT 图文全量提取完成！")
                except ImportError:
                    ui_container.error("❌ 缺少 pptx 库，请执行 `pip install python-pptx`")
                    continue
                    
            else:
                ui_container.warning(f"👁️ 正在激活视觉引擎，逐页提取新文件: {file.name} ...")
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                safe_filename = f"{timestamp}_{file.name}"
                file_path = os.path.join(UPLOAD_DIR, safe_filename)
                
                with open(file_path, "wb") as f: f.write(file_bytes)
                    
                image_paths = []
                if ext == '.pdf':
                    from core.parsers.document_engine import convert_pdf_to_images
                    image_paths.extend(convert_pdf_to_images(file_path, TEMP_IMG_DIR, max_pages))
                else:
                    image_paths.append(file_path)
                    
                if max_pages: image_paths = image_paths[:max_pages]
                total_pages = len(image_paths)
                if total_pages == 0:
                    ui_container.error(f"❌ {file.name} 提取失败！")
                    continue
                    
                progress_bar = ui_container.progress(0)
                for i, path in enumerate(image_paths):
                    res = process_single_page(path, i + 1)
                    all_content += f"\n\n> 📁 **[来源文件：{base_name}]** - 第 {i+1} 页提取内容\n{res}\n"
                    progress_bar.progress((i + 1) / total_pages)
                progress_bar.empty()
                
                try:
                    if os.path.exists(file_path): os.remove(file_path)
                    for img_path in image_paths:
                        if os.path.exists(img_path): os.remove(img_path)
                    ui_container.success(f"✅ {file.name} 解析完成，已存入物理指纹缓存池！")
                except Exception as e:
                    pass

            if all_content:
                with open(cache_file_path, "w", encoding="utf-8") as f: 
                    f.write(all_content)
                result_dict[base_name] = all_content
                
    return result_dict

tab1, tab2, tab3, tab4 = st.tabs([
    "🚀 单份文档智能解析", 
    "⚔️ 多公司竞品横评",
    "📈 历史纵向趋势推演",
    "📚 金牌模板库 (AI经验池)" 
])

# ---------------------------------------------------------
# 工作流 A：单份文档智能解析 (带对话式二次修改)
# ---------------------------------------------------------
with tab1:
    st.markdown("###### 📥 上传待解读文档 (支持混传 PDF / JPG / MD / Word / PPT)")
    uploaded_files = st.file_uploader("请拖拽文件至此", type=["pdf", "png", "jpg", "jpeg", "md", "docx", "pptx"], accept_multiple_files=True, key="tab1_uploader")
    
    col_req, col_tpl = st.columns([2, 1])
    with col_req:
        user_requirement_full = st.text_area(
            "🎯 自定义分析侧重点 (选填)",
            placeholder=(
                "不填则 AI 自动识别文档类型并自适应分析视角。\n"
                "也可手动指定侧重，例如：\n"
                "• 财报：重点分析毛利率与现金流质量\n"
                "• 政策文件：聚焦落地路径与关键时间节点"
            ),
            height=120
        )
    with col_tpl:
        options = ["🤖 AI 自动匹配金牌范例 (推荐)", "❌ 不参考金牌经验 (默认风格)"] + get_available_templates()
        selected_strategy = st.selectbox("🎯 选择报告行文风格：", options, key="tab1_strategy")
    
    col_btn1, col_btn2, col_btn3 = st.columns(3)
    with col_btn1: 
        btn_extract_only = st.button("📝 仅提取数据底稿", use_container_width=True)
    with col_btn2: 
        btn_quick_scan = st.button("⚡ 快速智能研判", type="primary", use_container_width=True) 
    with col_btn3: 
        btn_full_pipeline = st.button("🚀 红蓝军深度研判", use_container_width=True) 

    # 1. 触发生成逻辑
    if btn_extract_only or btn_quick_scan or btn_full_pipeline:
        if not uploaded_files:
            st.warning("⚠️ 请先上传文件！")
            st.stop()
            
        status_container = st.container()
        text_dict = parse_files_to_text_dict(uploaded_files, max_pages, status_container, use_cache)
        
        all_content = "\n".join(text_dict.values())
        base_name = os.path.splitext(uploaded_files[0].name)[0] if len(uploaded_files) == 1 else "多文件合并"
        
        # 缓存基础数据到 Session State，供后续对话参考
        st.session_state.tab1_docs = all_content
        st.session_state.tab1_base_name = base_name
        st.session_state.tab1_report_title = "深度洞察与业务研判报告" if btn_full_pipeline else "快速智能解读报告"
        
        # 模式 1：仅提取底稿
        if btn_extract_only:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            file_name_md = f"{base_name}_底稿_{timestamp}.md"
            temp_md_path = os.path.join(OUTPUT_DIR, file_name_md)
            with open(temp_md_path, "w", encoding="utf-8") as f: f.write(all_content)
            st.info("💡 提取完毕！您可以直接下载该 Markdown 文件用于预处理。")
            with open(temp_md_path, "r", encoding="utf-8") as f:
                st.download_button("⬇️ 一键下载 MD 数据底稿", data=f, file_name=file_name_md, mime="text/markdown", type="primary")
            st.stop()
            
        # 模式 2 & 3：研判逻辑
        analysis_mode = "deep" if btn_full_pipeline else "quick" 
        
        templates_str = ""
        ai_thinking_log = ""
        if btn_full_pipeline and selected_strategy != "❌ 不参考金牌经验 (默认风格)":
            with st.status("🧠 正在检索并研判人类金牌经验库...", expanded=True) as status:
                templates_str, ai_thinking_log = get_style_templates(all_content, selected_strategy, status)
            if selected_strategy.startswith("🤖") and ai_thinking_log:
                st.info(f"🤖 **大模型主编的选版笔记**：\n\n{ai_thinking_log}")

        style_instruction = GET_USER_PRIORITY(user_requirement_full) + UI_COGNITIVE_SINGLE
        if btn_full_pipeline and templates_str:
            style_instruction += GET_STYLE_FUSION(templates_str, report_type="single") + UI_CHART_MERMAID
        elif btn_quick_scan:
            style_instruction += "\n\n请直接基于文档进行结构化解读，无需执行红蓝军博弈推演。"

        with st.status('专家团队正在阅卷，请耐心等待...', expanded=True) as status_ui:
            summary, docx_path = generate_final_summary(
                all_content, 
                user_req=user_requirement_full, 
                style_instruction=style_instruction, 
                status_ui=status_ui,
                mode=analysis_mode 
            )
            status_ui.update(label="✅ 研判推演与物理封装完毕！", state="complete", expanded=False)
            
        st.session_state.tab1_report = summary
        st.session_state.tab1_docx_path = docx_path
        st.session_state.tab1_history = []
        st.rerun() 

    # 2. 渲染报告内容与对话框
    if st.session_state.tab1_report:
        st.success("🎉 研判报告已生成（当前为最新版本）！")
        st.markdown("---")
        st.markdown(st.session_state.tab1_report, unsafe_allow_html=True)
        
        final_md_content = (
            f"# AI {st.session_state.tab1_report_title}\n\n"
            f"{st.session_state.tab1_report}\n\n"
            f"---\n"
            f"## 📚 附录：原始底层数据\n"
            f"<details>\n"
            f"<summary>👉 点击展开查看各页 OCR 原始核心数据</summary>\n\n"
            f"{st.session_state.tab1_docs}\n"
            f"</details>"
        )
        
        render_export_buttons(final_md_content, st.session_state.tab1_base_name, "研判报告", st.session_state.tab1_docx_path)

        st.markdown("---")
        st.markdown("#### 💬 对话主编：持续修改与润色报告")
        for msg in st.session_state.tab1_history:
            with st.chat_message(msg["role"]):
                st.markdown(msg["content"])

        if prompt := st.chat_input("您可以指令：'删掉冗长部分'、'调整为激进语气' 或 '帮我加一段总结'", key="chat_tab1"):
            st.session_state.tab1_history.append({"role": "user", "content": prompt})
            with st.chat_message("user"): st.markdown(prompt)
            with st.chat_message("assistant"):
                # 👇 替换原来的 spinner
                st.markdown("🤖 **主编正在敲键盘修改中... (实时流式预览)**")
                stream_box = st.empty() # 创建一个空容器
                
                from modules.multi_compare.main import revise_report
                new_report = revise_report(
                    current_report=st.session_state.tab1_report, 
                    user_feedback=prompt, 
                    original_docs=st.session_state.tab1_docs,
                    status_ui=stream_box # 👈 把容器传给后端
                )
                
                st.session_state.tab1_report = new_report
                st.session_state.tab1_history.append({"role": "assistant", "content": "✅ 报告正文已按您的要求更新！您可以向上滚动查看最新内容。"})
                st.rerun()

# ---------------------------------------------------------
# 工作流 B：多公司竞品横评 (带对话式二次修改)
# ---------------------------------------------------------
with tab2:
    st.markdown("###### ⚔️ 上传多家公司的报告 (支持混传 PDF / JPG / MD / Word / PPT)")
    st.info("💡 提示：无论您传原文件还是MD底稿，请**直接用公司名称命名文件**。")
    
    col_req_b, col_tpl_b = st.columns([2, 1])
    with col_req_b:
        user_requirement_b = st.text_area(
            "🎯 自定义分析侧重点 (选填)",
            placeholder=(
                "不填则 AI 自动识别行业并切换竞品对比视角。\n"
                "也可手动指定，例如：\n"
                "• 重点对比各家在云服务市场上的战略差异\n"
                "• 聚焦研发投入强度与技术护城河的差距"
            ),
            height=120,
            key="tab2_req"
        )
    with col_tpl_b:
        options = ["🤖 AI 自动匹配金牌范例 (推荐)", "❌ 不参考金牌经验 (默认风格)"] + get_available_templates()
        selected_strategy_b = st.selectbox("🎯 选择报告行文风格：", options, key="tab2_strategy")

    compare_files = st.file_uploader("批量拖拽多个公司的文件至此", type=["pdf", "png", "jpg", "jpeg", "md", "docx", "pptx"], accept_multiple_files=True, key="tab2_uploader")
    
    if st.button("⚔️ 启动多 Agent 竞品横评", type="primary", key="btn_compare"):
        if not compare_files or len(compare_files) < 2:
            st.warning("⚠️ 进行横评至少需要上传 2 份不同公司的文件！")
            st.stop()
            
        status_container = st.container()
        company_dict = parse_files_to_text_dict(compare_files, max_pages, status_container, use_cache)
        
        # 缓存原始提取数据用于修改时的参考
        docs_str = "\n".join([f"【{k}】\n{v}" for k, v in company_dict.items() if k not in ["_STYLE_INSTRUCTION_", "_USER_REQ_"]])
        st.session_state.tab2_docs = docs_str
        names = [k for k in company_dict.keys() if k not in ["_STYLE_INSTRUCTION_", "_USER_REQ_"]]
        st.session_state.tab2_prefix = "vs".join(names[:3]) + ("等" if len(names)>3 else "")

        all_content_for_search = "\n".join(list(company_dict.values()))[:1000]
        
        templates_str = ""
        ai_thinking_log = ""
        if selected_strategy_b != "❌ 不参考金牌经验 (默认风格)":
            with st.status("🧠 正在检索并研判人类金牌经验库...", expanded=True) as status:
                templates_str, ai_thinking_log = get_style_templates(all_content_for_search, selected_strategy_b, status)
            if selected_strategy_b.startswith("🤖") and ai_thinking_log:
                st.info(f"🤖 **大模型主编的选版笔记**：\n\n{ai_thinking_log}")

        style_instruction = GET_USER_PRIORITY(user_requirement_b) + UI_COGNITIVE_COMPARE
        if templates_str:
            style_instruction += GET_STYLE_FUSION(templates_str, report_type="compare") + UI_CHART_MERMAID

        if style_instruction:
            company_dict["_STYLE_INSTRUCTION_"] = style_instruction
        if user_requirement_b.strip():
            company_dict["_USER_REQ_"] = user_requirement_b.strip()
            
        st.markdown("###### 🧠 多模态竞品大脑生成中")
        with st.status('正在强制大模型交叉校验数据，剔除无据猜测...', expanded=True) as status_ui_b:
            compare_summary = generate_compare_summary(company_dict, status_ui_b)
            status_ui_b.update(label="✅ 竞品横评数据处理与深度推演完毕！", state="complete", expanded=False)
            
        # 生成完毕，存入 Session State
        st.session_state.tab2_report = compare_summary
        st.session_state.tab2_history = []
        st.rerun()

    # 渲染报告与对话框
    if st.session_state.tab2_report:
        st.success("🎉 竞品横评报告已生成（当前为最新版本）！")
        st.markdown("---")
        st.markdown(st.session_state.tab2_report, unsafe_allow_html=True)
        
        final_md_content = (
            f"# ⚔️ 行业竞品横向对比研判报告\n\n"
            f"{st.session_state.tab2_report}\n\n"
            f"---\n"
            f"## 📚 附录：原始底层数据\n"
            f"<details>\n"
            f"<summary>👉 点击展开查看参与对比的各家原始数据</summary>\n\n"
            f"{st.session_state.tab2_docs}\n"
            f"</details>"
        )
        
        render_export_buttons(final_md_content, st.session_state.tab2_prefix, "竞品横评", None)

        st.markdown("---")
        st.markdown("#### 💬 对话主编：持续修改与润色横评报告")
        for msg in st.session_state.tab2_history:
            with st.chat_message(msg["role"]): st.markdown(msg["content"])

        if prompt := st.chat_input("您可以指令：'突出A公司在海外的优势' 或 '精简结论部分'", key="chat_tab2"):
            st.session_state.tab2_history.append({"role": "user", "content": prompt})
            with st.chat_message("user"): st.markdown(prompt)
            with st.chat_message("assistant"):
                # 👇 替换原来的 spinner
                st.markdown("🤖 **主编正在敲键盘修改横评报告中... (实时流式预览)**")
                stream_box = st.empty()
                
                from modules.multi_compare.main import revise_report
                new_report = revise_report(
                    current_report=st.session_state.tab2_report, 
                    user_feedback=prompt, 
                    original_docs=st.session_state.tab2_docs,
                    status_ui=stream_box # 👈 把容器传给后端
                )
                
                st.session_state.tab2_report = new_report
                st.session_state.tab2_history.append({"role": "assistant", "content": "✅ 横评报告正文已按您的要求更新！"})
                st.rerun()

# ---------------------------------------------------------
# 工作流 C：历史纵向趋势推演 (带对话式二次修改)
# ---------------------------------------------------------
with tab3:
    st.markdown("###### 📈 上传连续多年的报告 (支持混传 PDF / JPG / MD / Word / PPT)")
    st.info("💡 提示：无论您传原文件还是MD底稿，请**直接用年份命名文件**。")
    
    col_req_c, col_tpl_c = st.columns([2, 1])
    with col_req_c:
        user_requirement_c = st.text_area(
            "🎯 自定义分析侧重点 (选填)",
            placeholder=(
                "不填则 AI 自动识别行业并切换纵向演进分析视角。\n"
                "也可手动指定，例如：\n"
                "• 重点推演近三年研发投入与商业化成果的关系\n"
                "• 聚焦核心业务增速的拐点年份与战略归因"
            ),
            height=120,
            key="tab3_req"
        )
    with col_tpl_c:
        options = ["🤖 AI 自动匹配金牌范例 (推荐)", "❌ 不参考金牌经验 (默认风格)"] + get_available_templates()
        selected_strategy_c = st.selectbox("🎯 选择报告行文风格：", options, key="tab3_strategy")

    trend_files = st.file_uploader("批量拖拽多年的文件至此", type=["pdf", "png", "jpg", "jpeg", "md", "docx", "pptx"], accept_multiple_files=True, key="tab3_uploader")
    
    if st.button("📈 启动历史趋势推演", type="primary", key="btn_trend"):
        if not trend_files or len(trend_files) < 2:
            st.warning("⚠️ 进行趋势推演至少需要上传 2 个年份的文件！")
            st.stop()
            
        status_container = st.container()
        yearly_dict = parse_files_to_text_dict(trend_files, max_pages, status_container, use_cache)
        
        # 缓存原始提取数据用于修改时的参考
        docs_str = "\n".join([f"【{k}】\n{v}" for k, v in yearly_dict.items() if k not in ["_STYLE_INSTRUCTION_", "_USER_REQ_"]])
        st.session_state.tab3_docs = docs_str
        years = sorted([k for k in yearly_dict.keys() if k not in ["_STYLE_INSTRUCTION_", "_USER_REQ_"]])
        st.session_state.tab3_prefix = f"{years[0]}至{years[-1]}年" if len(years) > 1 else years[0]

        all_content_for_search = "\n".join(list(yearly_dict.values()))[:1000]
        
        templates_str = ""
        ai_thinking_log = ""
        if selected_strategy_c != "❌ 不参考金牌经验 (默认风格)":
            with st.status("🧠 正在检索并研判人类金牌经验库...", expanded=True) as status:
                templates_str, ai_thinking_log = get_style_templates(all_content_for_search, selected_strategy_c, status)
            if selected_strategy_c.startswith("🤖") and ai_thinking_log:
                st.info(f"🤖 **大模型主编的选版笔记**：\n\n{ai_thinking_log}")

        style_instruction = GET_USER_PRIORITY(user_requirement_c) + UI_COGNITIVE_TREND
        if templates_str:
            style_instruction += GET_STYLE_FUSION(templates_str, report_type="trend") + UI_CHART_MERMAID

        if style_instruction:
            yearly_dict["_STYLE_INSTRUCTION_"] = style_instruction
        if user_requirement_c.strip():
            yearly_dict["_USER_REQ_"] = user_requirement_c.strip()

        st.markdown("###### 🧠 历史趋势大脑生成中")
        with st.status('正在梳理历年时间轴，甄别断层数据...', expanded=True) as status_ui_c:
            trend_summary = generate_trend_summary(yearly_dict, status_ui_c)
            status_ui_c.update(label="✅ 历年趋势数据处理与深度推演完毕！", state="complete", expanded=False)
            
        # 生成完毕，存入 Session State
        st.session_state.tab3_report = trend_summary
        st.session_state.tab3_history = []
        st.rerun()

    # 渲染报告与对话框
    if st.session_state.tab3_report:
        st.success("🎉 纵向战略演进报告已生成（当前为最新版本）！")
        st.markdown("---")
        st.markdown(st.session_state.tab3_report, unsafe_allow_html=True)
        
        final_md_content = (
            f"# 📈 企业纵向战略演进与周期复盘报告\n\n"
            f"{st.session_state.tab3_report}\n\n"
            f"---\n"
            f"## 📚 附录：历年底层数据\n"
            f"<details>\n"
            f"<summary>👉 点击展开查看各年份原始底层数据</summary>\n\n"
            f"{st.session_state.tab3_docs}\n"
            f"</details>"
        )
        
        render_export_buttons(final_md_content, st.session_state.tab3_prefix, "演进趋势", None)

        st.markdown("---")
        st.markdown("#### 💬 对话主编：持续修改与润色趋势推演报告")
        for msg in st.session_state.tab3_history:
            with st.chat_message(msg["role"]): st.markdown(msg["content"])

        if prompt := st.chat_input("您可以指令：'重点归因2023年的下滑' 或 '加一段未来预测'", key="chat_tab3"):
            st.session_state.tab3_history.append({"role": "user", "content": prompt})
            with st.chat_message("user"): st.markdown(prompt)
            with st.chat_message("assistant"):
                # 👇 替换原来的 spinner
                st.markdown("🤖 **主编正在敲键盘修改趋势报告中... (实时流式预览)**")
                stream_box = st.empty()
                
                from modules.multi_compare.main import revise_report
                new_report = revise_report(
                    current_report=st.session_state.tab3_report, 
                    user_feedback=prompt, 
                    original_docs=st.session_state.tab3_docs,
                    status_ui=stream_box # 👈 把容器传给后端
                )
                
                st.session_state.tab3_report = new_report
                st.session_state.tab3_history.append({"role": "assistant", "content": "✅ 趋势报告正文已按您的要求更新！"})
                st.rerun()

# ---------------------------------------------------------
# Tab 4 - 经验库沉淀区 
# ---------------------------------------------------------
with tab4:
    st.markdown("#### 📥 投喂人工极品报告，让大模型越用越聪明")
    template_files = st.file_uploader("上传历史优秀报告作为金牌模板", accept_multiple_files=True, key="template_uploader")
    
    if st.button("🚀 提炼并入库经验池", type="secondary"):
        if not template_files:
            st.warning("请先上传文件！")
        else:
            status_container = st.container()
            with st.status("🧠 正在吸收人类智慧入库...", expanded=True) as status:
                embeddings = get_embeddings()
                db_path = os.path.join(TEMPLATE_DB_DIR, "index.faiss")
                vectorstore = FAISS.load_local(TEMPLATE_DB_DIR, embeddings, allow_dangerous_deserialization=True) if os.path.exists(db_path) else None
                
                status.write("正在调用全系统指纹高速缓存与分拣引擎...")
                tpl_dict = parse_files_to_text_dict(template_files, max_pages, status_container, enable_cache=use_cache)
                
                new_docs = []
                for base_name, content in tpl_dict.items():
                    md_filename = f"{base_name}.md"
                    with open(os.path.join(TEMPLATE_MD_DIR, md_filename), "w", encoding="utf-8") as f:
                        f.write(content)
                    
                    feature_text = content[:1500]
                    new_docs.append(Document(page_content=feature_text, metadata={"source": md_filename}))
                    status.write(f"✅ `{base_name}` 经验已成功向量化并吸收！")
                
                if new_docs:
                    if vectorstore is None: vectorstore = FAISS.from_documents(new_docs, embeddings)
                    else: vectorstore.add_documents(new_docs)
                    vectorstore.save_local(TEMPLATE_DB_DIR)
                    status.update(label="🎉 经验库更新完成！全局指纹系统已打通！", state="complete")
    
    st.divider()
    st.markdown("#### 🏆 已收录的金牌模板清单")
    existing_tpls = get_available_templates()
    if existing_tpls:
        for tpl in existing_tpls: st.markdown(f"- 📄 `{tpl}`")
    else: st.markdown("*当前经验库为空，请上传您的第一份金牌报告！*")